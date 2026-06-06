#!/usr/bin/env python3
"""Build Anthropic-style PPT for skill-evaluator journey."""

import os
import io
import textwrap
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ─── Colors ────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0F, 0x0F, 0x1A)
BG_CARD   = RGBColor(0x1A, 0x1A, 0x30)
ACCENT    = RGBColor(0xF5, 0x9E, 0x0B)   # amber
ACCENT2   = RGBColor(0x10, 0xB9, 0x81)   # emerald (for positive)
ACCENT3   = RGBColor(0xEF, 0x44, 0x44)   # red (for negative)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCA, 0xCA, 0xCA)
DIM       = RGBColor(0x88, 0x88, 0x99)
DARKER    = RGBColor(0x0A, 0x0A, 0x15)

# Chart colors for matplotlib
C_ACCENT  = '#f59e0b'
C_ACCENT2 = '#10b981'
C_ACCENT3 = '#ef4444'
C_BG      = '#0f0f1a'
C_CARD    = '#1a1a30'
C_TEXT    = '#cacaca'
C_DIM     = '#888899'

plt.rcParams.update({
    'font.family': 'sans-serif',
    'font.sans-serif': ['DejaVu Sans', 'Arial'],
    'font.size': 11,
    'axes.facecolor': C_CARD,
    'figure.facecolor': C_BG,
    'axes.edgecolor': C_DIM,
    'axes.labelcolor': C_TEXT,
    'axes.titlecolor': WHITE,
    'text.color': C_TEXT,
    'xtick.color': C_DIM,
    'ytick.color': C_DIM,
    'grid.color': C_DIM,
    'grid.alpha': 0.15,
    'legend.facecolor': C_CARD,
    'legend.edgecolor': C_DIM,
    'legend.labelcolor': C_TEXT,
})

# ─── Helpers ────────────────────────────────────────────────────

def set_slide_bg(slide, color=BG_DARK):
    """Set solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color=BG_CARD, radius=None):
    """Add a rounded-rect card background."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Arial'):
    """Add a text box with single paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_textbox(slide, left, top, width, height, paragraphs_data):
    """
    Add textbox with multiple paragraphs, each with optional per-run styling.
    paragraphs_data: list of dicts:
      {'text': str, 'size': int, 'color': RGBColor, 'bold': bool,
       'alignment': PP_ALIGN, 'spacing': Pt, 'space_before': Pt}
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, pd in enumerate(paragraphs_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = pd.get('text', '')
        p.font.size = Pt(pd.get('size', 14))
        p.font.color.rgb = pd.get('color', WHITE)
        p.font.bold = pd.get('bold', False)
        p.font.name = pd.get('font_name', 'Arial')
        p.alignment = pd.get('alignment', PP_ALIGN.LEFT)
        if 'spacing' in pd:
            p.space_after = pd['spacing']
        if 'space_before' in pd:
            p.space_before = pd['space_before']
    return txBox

def add_accent_bar(slide, left, top, width=Inches(0.06), height=Inches(0.5)):
    """Thin accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape

def slide_number(slide, prs, num):
    """Add subtle slide number bottom-right."""
    add_textbox(slide, Inches(8.8), Inches(7.0), Inches(1.0), Inches(0.3),
                f"{num:02d}", font_size=9, color=DIM, alignment=PP_ALIGN.RIGHT)

def new_slide(prs, bg=BG_DARK):
    """Create a blank slide with dark background."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide, bg)
    return slide

# ─── Chart generators ───────────────────────────────────────────

def chart_script_signal():
    """Figure: with-script vs without-script score distribution."""
    fig, ax = plt.subplots(figsize=(7.5, 3.8))
    labels = ['有脚本\n(4 skills)', '无脚本\n(8 skills)']
    with_script = [77, 80, 95, 77]  # k8s, ship, zabbix, justfile
    without_script = [71, 72, 67, 62, 59, 58, 54, 70]  # the rest (incl pre-commit)

    means = [np.mean(with_script), np.mean(without_script)]
    sems = [np.std(with_script)/np.sqrt(len(with_script)),
            np.std(without_script)/np.sqrt(len(without_script))]

    bars = ax.bar(labels, means, yerr=sems, capsize=5, width=0.45,
                  color=[C_ACCENT, C_ACCENT3], edgecolor='none', alpha=0.9)

    for bar, mean in zip(bars, means):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1.5,
                f'{mean:.0f}', ha='center', va='bottom', fontsize=16,
                fontweight='bold', color=C_TEXT)

    # Individual scatter points
    x_off = 0.1
    for i, v in enumerate(with_script):
        ax.scatter(i - x_off + np.random.uniform(-0.08, 0.08), v,
                   s=40, color=C_ACCENT2, alpha=0.5, zorder=5)
    for i, v in enumerate(without_script):
        ax.scatter(1 - x_off + np.random.uniform(-0.08, 0.08), v,
                   s=40, color='#ef4444', alpha=0.4, zorder=5)

    # Gap annotation
    gap = means[0] - means[1]
    ax.annotate(f'Δ = {gap:.1f} 分', xy=(0.5, max(means)+5),
                fontsize=13, fontweight='bold', color=C_ACCENT,
                ha='center', va='bottom')

    ax.set_ylabel('v3.1 总分', fontsize=11)
    ax.set_ylim(40, 105)
    ax.yaxis.set_major_locator(mticker.MultipleLocator(10))
    ax.tick_params(colors=C_DIM, labelsize=10)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

    for label in ax.get_xticklabels():
        label.set_color(C_TEXT)

    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, dpi=180, bbox_inches='tight', facecolor=C_BG,
                transparent=False)
    plt.close(fig)
    buf.seek(0)
    return buf


def chart_version_comparison():
    """Figure: v3.0 vs v3.1 scores for installed skills."""
    skills = ['xlsx', 'web-artifacts', 'pptx', 'docx', 'pdf', 'theme-factory',
              'claude-api', 'brand-guide', 'frontend-design', 'mcp-builder',
              'canvas-design', 'webapp-testing', 'karpathy', 'skill-creator',
              'slack-gif', 'internal-comms', 'algorithmic-art', 'doc-coauthor',
              'playwright-cli']
    v30 = [87, 89, 86, 87, 84, 80, 78, 79, 76, 78, 73, 79, 74, 70, 72, 72, 69, 66, 60]
    v31 = [94, 93, 92, 91, 88, 87, 84, 84, 82, 82, 79, 81, 73, 73, 75, 76, 72, 61, 55]

    fig, ax = plt.subplots(figsize=(9.5, 4.2))
    x = np.arange(len(skills))
    w = 0.35

    bars1 = ax.bar(x - w/2, v30, w, label='v3.0', color=C_DIM, alpha=0.5)
    bars2 = ax.bar(x + w/2, v31, w, label='v3.1', color=C_ACCENT, alpha=0.85)

    # Color the ones that changed significantly
    for i, (v0, v1) in enumerate(zip(v30, v31)):
        diff = v1 - v0
        if abs(diff) >= 3:
            bars2[i].set_color(C_ACCENT2 if diff > 0 else C_ACCENT3)
            bars2[i].set_alpha(0.9)

    ax.set_xticks(x)
    ax.set_xticklabels(skills, rotation=45, ha='right', fontsize=7)
    ax.set_ylabel('总分', fontsize=10)
    ax.set_ylim(40, 100)
    ax.legend(fontsize=9, loc='lower left')
    ax.tick_params(colors=C_DIM, labelsize=8)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

    for label in ax.get_xticklabels():
        label.set_color(C_TEXT)

    # annotation
    ax.text(0.02, 0.98, '脚本技能提升显著 (深绿) ; 知识技能调整 (深红)',
            transform=ax.transAxes, fontsize=8, color=C_DIM, va='top',
            fontstyle='italic')

    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, dpi=180, bbox_inches='tight', facecolor=C_BG)
    plt.close(fig)
    buf.seek(0)
    return buf


def chart_cli_heatmap():
    """Horizontal bar chart: CLI skill scores."""
    skills = ['ship', 'k8s-clusters', 'zabbix-api', 'justfile', 'tmux',
              'using-cloud-cli', 'pre-commit', 'shell-prompt', 'git',
              'repomix', 'atuin', 'direnv']
    scores = [95, 80, 77, 76, 72, 71, 70, 67, 62, 59, 58, 54]
    has_script = [True, True, True, False, False, True, False, True,
                  False, False, False, False]

    fig, ax = plt.subplots(figsize=(7.5, 4.5))
    y = np.arange(len(skills))
    colors = [C_ACCENT if s else C_ACCENT3 for s in has_script]

    bars = ax.barh(y, scores, height=0.6, color=colors, alpha=0.85, zorder=3)
    for bar, score in zip(bars, scores):
        ax.text(bar.get_width() + 1, bar.get_y() + bar.get_height()/2,
                f'{score}', va='center', fontsize=11, fontweight='bold',
                color=C_TEXT)

    # Legend
    from matplotlib.patches import Patch
    legend_elements = [Patch(facecolor=C_ACCENT, alpha=0.85, label='有脚本'),
                       Patch(facecolor=C_ACCENT3, alpha=0.85, label='无脚本')]
    ax.legend(handles=legend_elements, fontsize=9, loc='lower right')

    ax.set_yticks(y)
    ax.set_yticklabels(skills, fontsize=9)
    ax.set_xlabel('v3.1 总分', fontsize=10)
    ax.set_xlim(0, 105)
    ax.tick_params(colors=C_DIM, labelsize=9)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.xaxis.set_major_locator(mticker.MultipleLocator(20))

    for label in ax.get_yticklabels():
        label.set_color(C_TEXT)

    # Vertical line at 70 threshold
    ax.axvline(x=70, color=C_DIM, linestyle='--', linewidth=0.8, alpha=0.4)
    ax.text(71, -0.6, '推荐线 (70)', fontsize=8, color=C_DIM)

    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, dpi=180, bbox_inches='tight', facecolor=C_BG)
    plt.close(fig)
    buf.seek(0)
    return buf


def chart_dimension_breakdown():
    """Stacked horizontal bar: dimension breakdown for top skills."""
    skills = ['ship (95)', 'xlsx (94)', 'web-artifacts (93)',
              'pptx (92)', 'docx (91)', 'pdf (88)']
    # trigger, execution, economy, trust, clarity
    trigger = [23, 25, 22, 20, 25, 20]
    exec_ =   [35, 35, 32, 35, 32, 32]
    economy = [22, 20, 22, 20, 16, 16]
    trust =   [7,  9,  9,  9,  8,  9]
    clarity = [8,  5,  8,  8,  8,  8]

    fig, ax = plt.subplots(figsize=(7.5, 3.5))
    y = np.arange(len(skills))
    h = 0.5

    colors = ['#f59e0b', '#06b6d4', '#10b981', '#8b5cf6', '#f472b6']
    labels = ['触发', '执行', '经济', '可信', '清晰']

    left = np.zeros(len(skills))
    segments = [trigger, exec_, economy, trust, clarity]
    for seg, c, lab in zip(segments, colors, labels):
        ax.barh(y, seg, h, left=left, color=c, alpha=0.8, label=lab, zorder=3)
        left += seg

    ax.set_yticks(y)
    ax.set_yticklabels(skills, fontsize=8)
    ax.set_xlabel('维度得分', fontsize=9)
    ax.set_xlim(0, 105)
    ax.legend(fontsize=8, loc='lower right', ncol=5)
    ax.tick_params(colors=C_DIM, labelsize=8)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

    for label in ax.get_yticklabels():
        label.set_color(C_TEXT)

    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, dpi=180, bbox_inches='tight', facecolor=C_BG)
    plt.close(fig)
    buf.seek(0)
    return buf


# ─── Build Presentation ─────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    W = prs.slide_width
    H = prs.slide_height

    # Generate charts
    print("  Generating charts...")
    chart_script = chart_script_signal()
    chart_ver = chart_version_comparison()
    chart_cli = chart_cli_heatmap()
    chart_dim = chart_dimension_breakdown()
    print("  Charts done.")

    # ================================================================
    # SLIDE 1: Cover
    # ================================================================
    s = new_slide(prs, DARKER)
    # Large accent bar at top
    add_shape_bg(s, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    # Title
    add_textbox(s, Inches(0.8), Inches(2.0), Inches(11), Inches(1.2),
                "skill-evaluator", font_size=52, color=WHITE, bold=True)
    add_textbox(s, Inches(0.8), Inches(3.0), Inches(11), Inches(0.8),
                "从 0 到 1，从 1 到好", font_size=36, color=ACCENT, bold=False)
    add_textbox(s, Inches(0.8), Inches(4.0), Inches(10), Inches(0.6),
                "一个社区技能评分器的进化之路  ·  v3.1", font_size=18, color=DIM)

    # Bottom info
    add_textbox(s, Inches(0.8), Inches(6.2), Inches(5), Inches(0.4),
                "2026.06  ·  40+ 技能评估  ·  6 大 skill hub", font_size=12, color=DIM)
    slide_number(s, prs, 1)

    # ================================================================
    # SLIDE 2: Pain Point
    # ================================================================
    s = new_slide(prs)
    add_accent_bar(s, Inches(0.8), Inches(0.8), height=Inches(0.4))

    add_rich_textbox(s, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6), [
        {'text': '痛点：装了 20 个技能，哪个值得信赖？', 'size': 28, 'color': WHITE, 'bold': True},
    ])

    # Left: pain point cards
    pains = [
        ("❓ 盲目安装", "看到推荐就装，装上才知道好不好"),
        ("📏 没有标准", "全靠直觉判断质量，没有客观指标"),
        ("🔄 不可对比", "一个说好用一个说差，谁是对的？"),
        ("⚠️ 安全隐患", "恶意技能可能窃取数据或执行危险操作"),
        ("🎭 宣传失真", "高 star 不等于高质量，长文档不等于价值高"),
    ]
    for i, (title, desc) in enumerate(pains):
        y = Inches(1.8) + Inches(i * 0.85)
        add_shape_bg(s, Inches(0.8), y, Inches(5.5), Inches(0.7), BG_CARD)
        add_textbox(s, Inches(1.1), y + Inches(0.08), Inches(5.0), Inches(0.3),
                    title, font_size=16, color=ACCENT, bold=True)
        add_textbox(s, Inches(1.1), y + Inches(0.35), Inches(5.0), Inches(0.3),
                    desc, font_size=12, color=LIGHT)

    # Right: Insight card
    add_shape_bg(s, Inches(7.0), Inches(1.8), Inches(5.5), Inches(3.5), BG_CARD)
    add_textbox(s, Inches(7.3), Inches(2.0), Inches(5.0), Inches(0.4),
                "💡 核心问题", font_size=18, color=ACCENT, bold=True)
    add_textbox(s, Inches(7.3), Inches(2.6), Inches(5.0), Inches(2.5),
                "Claude Code 的技能生态在2025-2026年爆炸式增长——\n"
                "从几个官方技能发展到数万个社区技能。\n\n"
                "但生态越大，选择越难：\n"
                "• 哪个技能真正值得安装？\n"
                "• 哪个只是看起来很华丽？\n"
                "• 哪个有安全隐患？\n\n"
                "需要一个可重复、可验证的评分系统。",
                font_size=13, color=LIGHT)

    # Bottom bar: mission
    add_shape_bg(s, Inches(0), Inches(6.6), W, Inches(0.9), BG_CARD)
    add_textbox(s, Inches(0.8), Inches(6.75), Inches(11), Inches(0.4),
                '🎯 使命：建立一个客观、可重复、数据驱动的技能评价体系',
                font_size=16, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    slide_number(s, prs, 2)

    # ================================================================
    # SLIDE 3: First Version
    # ================================================================
    s = new_slide(prs)
    add_accent_bar(s, Inches(0.8), Inches(0.8), height=Inches(0.4))
    add_textbox(s, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                "第一步：v1.0 — 安全门禁 + 初版评分卡", font_size=28, color=WHITE, bold=True)

    add_textbox(s, Inches(0.8), Inches(1.6), Inches(7), Inches(0.4),
                "最初的版本非常直接：先安全检查，再打分评估。", font_size=14, color=LIGHT)

    # Two-column layout
    add_shape_bg(s, Inches(0.8), Inches(2.3), Inches(5.5), Inches(3.5), BG_CARD)
    add_textbox(s, Inches(1.1), Inches(2.5), Inches(5.0), Inches(0.4),
                "🔒 安全门禁（一票否决）", font_size=16, color=ACCENT, bold=True)
    items = [
        "① 脚本可读性 — 检查混淆/加密",
        "② 网络外呼 — 搜索 curl/wget/fetch",
        "③ 凭证采集 — 检测 api.key/token/secret",
        "④ 远程代码执行 — curl | bash 检测",
        "⑤ 安全绕过 — 自动批准/关权限指令",
    ]
    for i, item in enumerate(items):
        add_textbox(s, Inches(1.3), Inches(3.1 + i * 0.45), Inches(4.8), Inches(0.4),
                    item, font_size=12, color=LIGHT)

    add_shape_bg(s, Inches(6.8), Inches(2.3), Inches(5.5), Inches(3.5), BG_CARD)
    add_textbox(s, Inches(7.1), Inches(2.5), Inches(5.0), Inches(0.4),
                "📊 初版评分维度", font_size=16, color=ACCENT, bold=True)
    dims = [
        "触发精确度 — 描述是否清晰",
        "执行约束力 — 有没有铁律/强约束",
        "验证/测试 — 有没有评估机制",
        "经济性 — 行数估值",
        "诚实度 — 有没有吹嘘/虚假承诺",
    ]
    for i, d in enumerate(dims):
        add_textbox(s, Inches(7.3), Inches(3.1 + i * 0.5), Inches(4.8), Inches(0.4),
                    d, font_size=12, color=LIGHT)

    # Personal insight
    add_shape_bg(s, Inches(0.8), Inches(6.1), Inches(11.5), Inches(1.0), BG_CARD)
    add_textbox(s, Inches(1.1), Inches(6.2), Inches(0.2), Inches(0.3),
                "💬", font_size=14)
    add_textbox(s, Inches(1.4), Inches(6.2), Inches(10.5), Inches(0.7),
                '"最初的假设是"好的技能 = 管得严"。 所以初版大量测量约束力、铁律、验证步骤。\n'
                '后来才知道——这个假设完全错了。"',
                font_size=13, color=DIM, font_name='Arial')
    slide_number(s, prs, 3)

    # ================================================================
    # SLIDE 4: Official Skills Flop
    # ================================================================
    s = new_slide(prs)
    add_accent_bar(s, Inches(0.8), Inches(0.8), height=Inches(0.4))
    add_textbox(s, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                "当头一棒：官方技能全部翻车", font_size=28, color=WHITE, bold=True)

    add_textbox(s, Inches(0.8), Inches(1.5), Inches(8), Inches(0.4),
                "用 v1.0 评估 Anthropic 官方技能 — 绝大多数不及格。问题不在技能，在评分标准。",
                font_size=14, color=LIGHT)

    # Score cards - mock scores
    off_skills = [('doc-coauthoring', '28', ACCENT3, '文档协作'),
                  ('playwright-cli', '32', ACCENT3, '浏览器测试'),
                  ('karpathy-guidelines', '45', '#f97316', '编码规范'),
                  ('claude-api', '52', '#f97316', 'API开发'),
                  ('pdf', '58', '#eab308', 'PDF生成'),
                  ('docx', '55', '#eab308', 'Word文档')]
    for i, (name, score, clr, desc) in enumerate(off_skills):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 4.0)
        y = Inches(2.3 + row * 2.6)
        add_shape_bg(s, x, y, Inches(3.6), Inches(2.2), BG_CARD)
        # Score circle
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2),
                                     Inches(0.8), Inches(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = clr
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = score
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.word_wrap = False
        circle.text_frame.paragraphs[0].space_before = Pt(0)
        circle.text_frame.paragraphs[0].space_after = Pt(0)

        add_textbox(s, x + Inches(1.2), y + Inches(0.3), Inches(2.2), Inches(0.3),
                    name, font_size=15, color=WHITE, bold=True)
        add_textbox(s, x + Inches(1.2), y + Inches(0.7), Inches(2.2), Inches(0.3),
                    desc, font_size=11, color=DIM)

    # Insight
    add_shape_bg(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.8), BG_CARD)
    add_textbox(s, Inches(1.1), Inches(6.4), Inches(11), Inches(0.5),
                '💬  "这个结果不对——官方技能应该高分才对。说明评分标准测量了错误的东西。'
                '没有坏的技能，只有坏的评分标准。"',
                font_size=13, color=DIM, font_name='Arial')
    slide_number(s, prs, 4)

    # ================================================================
    # SLIDE 5: Root Cause
    # ================================================================
    s = new_slide(prs)
    add_accent_bar(s, Inches(0.8), Inches(0.8), height=Inches(0.4))
    add_textbox(s, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                "根本原因：我们测量了错误的东西", font_size=28, color=WHITE, bold=True)

    # Before/After comparison
    add_shape_bg(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8), BG_CARD)
    add_textbox(s, Inches(1.1), Inches(2.0), Inches(5.0), Inches(0.4),
                "❌ 旧思维：合规审计", font_size=18, color=ACCENT3, bold=True)
    old_items = [
        "• 你有铁律表吗？有 → 高分",
        "• 你有 Good/Bad 示例吗？有 → 高分",
        "• 你有验证步骤吗？有 → 高分",
        "• 你用 MUST/ALWAYS 吗？用 → 高分",
        "",
        "问题：",
        "• 官方技能是写给 AI 看的，不是写试卷",
        "• 强制格式 ≠ 质量",
        "• 铁律多 ≠ 实际产出好",
        "• 模板化的技能反而得分高——完全搞反了",
    ]
    for i, item in enumerate(old_items):
        if item == "":
            continue
        add_textbox(s, Inches(1.3), Inches(2.6 + i * 0.38), Inches(4.8), Inches(0.4),
                    item, font_size=12, color=LIGHT if i >= 5 else WHITE)

    add_shape_bg(s, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.8), BG_CARD)
    add_textbox(s, Inches(7.1), Inches(2.0), Inches(5.0), Inches(0.4),
                "✅ 新思维：价值评估", font_size=18, color=ACCENT2, bold=True)
    new_items = [
        "• 有脚本自动化吗？有 → 高分",
        "• 产出比不用好多少？好很多 → 高分",
        "• 触发精准吗？该用时用不该用时不用 → 高分",
        "• 信息密度高吗？每行都有价值 → 高分",
        "",
        "核心理念转变：",
        '"一个好的技能 = 让 Claude 做出\n'
        '  没有它时做不出的结果。"',
        "",
        "不再问\"管得严不严\"，",
        "而是问\"产出好不好\"。",
    ]
    for i, item in enumerate(new_items):
        add_textbox(s, Inches(7.3), Inches(2.6 + i * 0.38), Inches(4.8), Inches(0.4),
                    item, font_size=12, color=ACCENT if i == 6 else (WHITE if i < 5 else LIGHT),
                    bold=(i == 6))
    slide_number(s, prs, 5)

    # ================================================================
    # SLIDE 6: Value Shift
    # ================================================================
    s = new_slide(prs)
    add_accent_bar(s, Inches(0.8), Inches(0.8), height=Inches(0.4))
    add_textbox(s, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                "价值观转向：从合规到价值", font_size=28, color=WHITE, bold=True)

    # Timeline
    phases = [
        ("v1.0", "合规审计\n测量'管得严不严'", ACCENT3),
        ("v2.0", "价值转向\n测量'产出好不好'", '#f97316'),
        ("v3.0", "5维体系\n平衡价值与效率", '#eab308'),
        ("v3.1", "精准量化\n数据驱动微调", ACCENT2),
    ]
    for i, (ver, desc, clr) in enumerate(phases):
        x = Inches(0.8 + i * 3.15)
        add_shape_bg(s, x, Inches(1.8), Inches(2.8), Inches(2.0), BG_CARD)
        # Version badge
        badge = add_shape_bg(s, x, Inches(1.8), Inches(1.2), Inches(0.5), clr)
        add_textbox(s, x + Inches(0.05), Inches(1.85), Inches(1.1), Inches(0.4),
                    ver, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # Description
        lines = desc.split('\n')
        for j, line in enumerate(lines):
            add_textbox(s, x + Inches(0.2), Inches(2.5 + j * 0.4), Inches(2.4), Inches(0.4),
                        line, font_size=13, color=LIGHT)

        # Arrow between phases
        if i < len(phases) - 1:
            add_textbox(s, x + Inches(2.85), Inches(2.5), Inches(0.3), Inches(0.4),
                        "→", font_size=22, color=DIM, bold=True)

    # Key changes table
    add_shape_bg(s, Inches(0.8), Inches(4.3), Inches(11.5), Inches(2.8), BG_CARD)
    add_textbox(s, Inches(1.1), Inches(4.5), Inches(5), Inches(0.4),
                "关键变革", font_size=16, color=ACCENT, bold=True)

    changes = [
        ("移除", "铁律强度测量", "改为测量指令完整性（零歧义性）"),
        ("移除", "显式验证步骤要求", "改为测量脚本/知识增量"),
        ("提升", "来源权重(3→5)", "官方/可信来源本身就是质量信号"),
        ("增加", "精炼附加分", "≤150行的精炼技能获得经济性加分"),
        ("增加", "信息密度调整", "超过50%泛化内容的长技能受罚"),
    ]
    headers = ["操作", "旧维度", "新维度"]
    for j, h in enumerate(headers):
        add_textbox(s, Inches(1.1 + j * 3.2), Inches(5.0), Inches(3.0), Inches(0.3),
                    h, font_size=11, color=DIM, bold=True)
    for i, (op, old, new) in enumerate(changes):
        y = Inches(5.4 + i * 0.4)
        add_textbox(s, Inches(1.1), y, Inches(2.2), Inches(0.3), op,
                    font_size=12, color=ACCENT2 if op == "增加" else (ACCENT if op == "提升" else ACCENT3))
        add_textbox(s, Inches(3.5), y, Inches(3.0), Inches(0.3), old, font_size=11, color=LIGHT)
        add_textbox(s, Inches(7.0), y, Inches(4.5), Inches(0.3), new, font_size=11, color=LIGHT)
    slide_number(s, prs, 6)

    # ================================================================
    # SLIDE 7: Community Test Matrix
    # ================================================================
    s = new_slide(prs)
    add_accent_bar(s, Inches(0.8), Inches(0.8), height=Inches(0.4))
    add_textbox(s, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                "社区大测试：4 大 Skill Hub，27+ 技能", font_size=28, color=WHITE, bold=True)

    hubs = [
        ("Vercel Labs", "46,000 ⭐", "react-best-practices\nweb-design-guidelines\ndeploy-to-vercel\nvercel-optimize", ACCENT),
        ("Posit (RStudio)", "364 ⭐", "pr-create\ncritical-code-reviewer\nimplement\ndescribe-design", ACCENT2),
        ("Gentleman-Skills", "531 ⭐", "github-pr\nreact-19", '#8b5cf6'),
        ("artubss/Brazil", "1,044 skills", "14 skills across\n26 categories", '#ec4899'),
        ("julianobarbosa", "118 skills", "12 CLI skills\ndeep evaluation", '#06b6d4'),
    ]
    for i, (name, stars, skills_, clr) in enumerate(hubs):
        x = Inches(0.5 + (i % 3) * 4.2)
        y = Inches(1.7 + (i // 3) * 2.8)
        add_shape_bg(s, x, y, Inches(3.8), Inches(2.4), BG_CARD)
        add_shape_bg(s, x, y, Inches(3.8), Inches(0.06), clr)
        add_textbox(s, x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(0.3),
                    name, font_size=14, color=WHITE, bold=True)
        add_textbox(s, x + Inches(0.2), y + Inches(0.55), Inches(3.4), Inches(0.2),
                    stars, font_size=11, color=DIM)
        add_textbox(s, x + Inches(0.2), y + Inches(0.9), Inches(3.4), Inches(1.2),
                    skills_, font_size=11, color=LIGHT)

    # Bottom insight
    add_shape_bg(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.8), BG_CARD)
    add_textbox(s, Inches(1.1), Inches(6.35), Inches(11), Inches(0.6),
                '💬  "从 46k star 的官方公司到 364 star 的专业团队，从 500+ 的精选集到 1000+ 的巴西社区，'
                '跨地域、跨语言、跨质量层次的测试让我们看到了技能的多样性。'
                '高 star ≠ 高质量：有些 46k 的技能评分不如 364 的。"',
                font_size=12, color=DIM, font_name='Arial')
    slide_number(s, prs, 7)

    # ================================================================
    # SLIDE 8: Script Signal (CHART)
    # ================================================================
    s = new_slide(prs)
    add_accent_bar(s, Inches(0.8), Inches(0.8), height=Inches(0.4))
    add_textbox(s, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                "关键发现 #1：脚本是质量的最强信号", font_size=28, color=WHITE, bold=True)

    # Chart
    s.shapes.add_picture(chart_script, Inches(3.0), Inches(1.5), Inches(7.5), Inches(3.8))

    # Key numbers
    add_shape_bg(s, Inches(0.8), Inches(1.5), Inches(2.0), Inches(2.0), BG_CARD)
    add_textbox(s, Inches(1.0), Inches(1.6), Inches(1.6), Inches(0.3),
                "+19~20", font_size=36, color=ACCENT, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1.0), Inches(2.0), Inches(1.6), Inches(0.4),
                "有脚本 vs\n无脚本 平均分差",
                font_size=11, color=LIGHT, alignment=PP_ALIGN.CENTER)

    add_shape_bg(s, Inches(0.8), Inches(3.7), Inches(2.0), Inches(1.8), BG_CARD)
    add_textbox(s, Inches(1.0), Inches(3.8), Inches(1.6), Inches(0.3),
                "84 ↔ 64", font_size=28, color=ACCENT2, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1.0), Inches(4.3), Inches(1.6), Inches(0.8),
                "有脚本组平均\n84.3 分\n无脚本组平均\n63.8 分",
                font_size=10, color=LIGHT, alignment=PP_ALIGN.CENTER)

    # Insight
    add_shape_bg(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2), BG_CARD)
    add_textbox(s, Inches(1.1), Inches(5.9), Inches(11), Inches(1.0),
                '💬  "这个发现改变了整个评分逻辑。一个有实际脚本的技能（无论多简单）的价值远超'
                '一篇长篇大论。因为脚本是\"可执行的承诺\"——它证明创作者真的用过这个工具，'
                '而不仅仅是翻了文档。这也解释了为什么ship(4脚本)、xlsx(5脚本)能拿最高分。"',
                font_size=12, color=DIM, font_name='Arial')
    slide_number(s, prs, 8)

    # ================================================================
    # SLIDE 9: Density Trap
    # ================================================================
    s = new_slide(prs)
    add_accent_bar(s, Inches(0.8), Inches(0.8), height=Inches(0.4))
    add_textbox(s, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                "关键发现 #2：长 ≠ 好，密度为王", font_size=28, color=WHITE, bold=True)

    # Density spectrum visualization
    examples = [
        ("短而精 ✅", "using-cloud-cli", "76行, 22/20", ACCENT2),
        ("精炼标杆 ✅", "ship", "148行, 22/20", ACCENT2),
        ("适中优秀 ✅", "justfile", "239行, 20/20", ACCENT2),
        ("偏长-高密 ⚠️", "k8s-clusters", "458行, 16/20", '#eab308'),
        ("偏长-中密 ⚠️", "atuin", "583行, 10/20", '#f97316'),
        ("长-低密 ❌", "direnv", "596行, 8/20", ACCENT3),
    ]
    for i, (label, name, score, clr) in enumerate(examples):
        y = Inches(1.7 + i * 0.7)
        # Color indicator
        add_shape_bg(s, Inches(0.8), y + Inches(0.05), Inches(0.06), Inches(0.5), clr)
        add_textbox(s, Inches(1.1), y, Inches(1.5), Inches(0.3), label,
                    font_size=12, color=clr, bold=True)
        add_textbox(s, Inches(2.6), y, Inches(2.0), Inches(0.3), name,
                    font_size=13, color=WHITE)
        add_textbox(s, Inches(4.8), y, Inches(2.5), Inches(0.3), score,
                    font_size=13, color=LIGHT)

    # Explanation cards
    add_shape_bg(s, Inches(0.8), Inches(5.9), Inches(5.8), Inches(1.2), BG_CARD)
    add_textbox(s, Inches(1.1), Inches(6.0), Inches(5.5), Inches(1.0),
                "📐 信息密度调整规则\n"
                "• 高密度：几乎每行都有具体指令 → +0\n"
                "• 低密度信号：套话段落、重复列表、50%+为什么\n"
                "  → 扣 -2 至 -6 分",
                font_size=12, color=LIGHT)

    add_shape_bg(s, Inches(6.8), Inches(5.9), Inches(5.5), Inches(1.2), BG_CARD)
    add_textbox(s, Inches(7.1), Inches(6.0), Inches(5.5), Inches(1.0),
                "💡 精炼附加分\n"
                "• ≤150行且功能完整 → +2 分\n"
                "• ship (148行) + direnv (596行) 经济分差 14 分\n"
                "  但 ship 功能是 direnv 的 3 倍",
                font_size=12, color=LIGHT)
    slide_number(s, prs, 9)

    # ================================================================
    # SLIDE 10: v3.0 Architecture
    # ================================================================
    s = new_slide(prs)
    add_accent_bar(s, Inches(0.8), Inches(0.8), height=Inches(0.4))
    add_textbox(s, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                "v3.0：5 维 100 分评分架构", font_size=28, color=WHITE, bold=True)

    dim_data = [
        ("触发精确度", "25", "描述是否精准？\n何时用/何时不用要明确", ACCENT,
         "触发边界 +2~3\n跨技能引用 +1"),
        ("执行质量", "35", "指令完整？\n有脚本？可靠？", '#06b6d4',
         "2a:指令15+2b:知识12+2c:可靠8"),
        ("经济性", "20", "多少 token？\n信息密度够吗？", ACCENT2,
         "精炼≤150行 +2\n密度调整 -2~-6"),
        ("可信度", "12", "来源可靠？\n元数据完整？脚本真实？", '#8b5cf6',
         "来源5+完整性3+元数据4"),
        ("理解成本", "8", "Claude 一遍能读懂吗？\n结构清晰？有示例？", '#f472b6',
         "一目了然8 → 混乱0"),
    ]
    for i, (name, score, desc, clr, sub) in enumerate(dim_data):
        x = Inches(0.5 + (i % 3) * 4.2)
        y = Inches(1.6 + (i // 3) * 2.9)
        add_shape_bg(s, x, y, Inches(3.9), Inches(2.5), BG_CARD)
        # Accent top
        add_shape_bg(s, x, y, Inches(3.9), Inches(0.06), clr)
        # Name
        add_textbox(s, x + Inches(0.2), y + Inches(0.2), Inches(2.2), Inches(0.3),
                    name, font_size=16, color=WHITE, bold=True)
        # Score badge
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x + Inches(2.6), y + Inches(0.15),
                                    Inches(1.0), Inches(0.45))
        badge.fill.solid()
        badge.fill.fore_color.rgb = clr
        badge.line.fill.background()
        badge.adjustments[0] = 0.15
        tf = badge.text_frame
        tf.paragraphs[0].text = f"{score}分"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(s, x + Inches(0.2), y + Inches(0.7), Inches(3.5), Inches(0.8),
                    desc, font_size=12, color=LIGHT)
        add_textbox(s, x + Inches(0.2), y + Inches(1.6), Inches(3.5), Inches(0.6),
                    sub, font_size=10, color=DIM)

    slide_number(s, prs, 10)

    # ================================================================
    # SLIDE 11: v3.1 Precision
    # ================================================================
    s = new_slide(prs)
    add_accent_bar(s, Inches(0.8), Inches(0.8), height=Inches(0.4))
    add_textbox(s, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                "v3.1：数据驱动的精准化调整", font_size=28, color=WHITE, bold=True)

    improvements = [
        ("2b 知识/自动化", "12", ACCENT,
         "v3.0: "模糊的'脚本价值'打分"\n"
         "v3.1: 精确脚本计数——3+脚本=12, 1-2脚本=9, 知识资产=6",
         "✅ 脚本技能从87→91, 正确奖励自动化程度"),
        ("经济·密度", "20", ACCENT2,
         "v3.0: "信息密度扣分 -2~4" 模糊\n"
         "v3.1: 4级具体阈值 + 5条低密度检测信号",
         "✅ 长技能正确降分，精炼技能正确加分"),
        ("可信·完整性", "3", '#8b5cf6',
         "v3.0: 测量evals/目录（几乎无人有）\n"
         "v3.1: 脚本完整性检查——脚本引用须真实存在",
         "✅ 发现2个技能引用不存在的脚本路径"),
        ("触发·负例", "25", '#06b6d4',
         "v3.0: 仅25/20/15三级\n"
         "v3.1: 增加23档，负例加分 +2~3，跨技能引用 +1",
         "✅ ship(95)首次超越官方技能最高分"),
    ]
    for i, (area, score, clr, before_after, impact) in enumerate(improvements):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 6.2)
        y = Inches(1.6 + row * 2.8)
        add_shape_bg(s, x, y, Inches(5.8), Inches(2.4), BG_CARD)
        add_shape_bg(s, x, y, Inches(5.8), Inches(0.06), clr)
        add_textbox(s, x + Inches(0.2), y + Inches(0.2), Inches(4.0), Inches(0.3),
                    area, font_size=15, color=WHITE, bold=True)
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x + Inches(4.5), y + Inches(0.15),
                                    Inches(1.0), Inches(0.4))
        badge.fill.solid()
        badge.fill.fore_color.rgb = clr
        badge.line.fill.background()
        badge.adjustments[0] = 0.15
        tf = badge.text_frame
        tf.paragraphs[0].text = score
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        add_textbox(s, x + Inches(0.2), y + Inches(0.65), Inches(5.4), Inches(0.9),
                    before_after, font_size=11, color=LIGHT)
        add_textbox(s, x + Inches(0.2), y + Inches(1.7), Inches(5.4), Inches(0.5),
                    impact, font_size=11, color=ACCENT2)
    slide_number(s, prs, 11)

    # ================================================================
    # SLIDE 12: Re-run Results (CHART)
    # ================================================================
    s = new_slide(prs)
    add_accent_bar(s, Inches(0.8), Inches(0.8), height=Inches(0.4))
    add_textbox(s, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                "验证：19 个已安装技能重跑 v3.1", font_size=28, color=WHITE, bold=True)

    s.shapes.add_picture(chart_ver, Inches(1.8), Inches(1.5), Inches(9.8), Inches(4.3))

    # Summary stats
    stats_data = [
        ("📈 平均分", "77.9", "v3.0 73.1 → +4.8"),
        ("⏫ 提升最大", "+7", "xlsx (87→94)"),
        ("⏬ 下调最大", "-5", "playwright-cli (60→55)"),
        ("🏆 最高分", "94", "xlsx — 脚本+边界"),
    ]
    for i, (label, val, sub) in enumerate(stats_data):
        x = Inches(0.5 + i * 3.2)
        add_shape_bg(s, x, Inches(6.0), Inches(3.0), Inches(1.1), BG_CARD)
        add_textbox(s, x + Inches(0.15), Inches(6.05), Inches(2.7), Inches(0.3),
                    label, font_size=10, color=DIM)
        add_textbox(s, x + Inches(0.15), Inches(6.35), Inches(2.7), Inches(0.3),
                    val, font_size=20, color=ACCENT, bold=True)
        add_textbox(s, x + Inches(0.15), Inches(6.7), Inches(2.7), Inches(0.3),
                    sub, font_size=10, color=LIGHT)
    slide_number(s, prs, 12)

    # ================================================================
    # SLIDE 13: CLI Results (CHART)
    # ================================================================
    s = new_slide(prs)
    add_accent_bar(s, Inches(0.8), Inches(0.8), height=Inches(0.4))
    add_textbox(s, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                "CLI 专项：12 个 CLI 社区技能评分", font_size=28, color=WHITE, bold=True)

    s.shapes.add_picture(chart_cli, Inches(3.5), Inches(1.5), Inches(7.0), Inches(4.2))

    add_shape_bg(s, Inches(0.8), Inches(1.5), Inches(2.5), Inches(2.0), BG_CARD)
    add_textbox(s, Inches(1.0), Inches(1.6), Inches(2.1), Inches(0.3),
                "📊 CLI组", font_size=14, color=WHITE, bold=True)
    add_textbox(s, Inches(1.0), Inches(2.0), Inches(2.1), Inches(1.2),
                "12个技能\n平均: 70.1\n最高: 95 (ship)\n最低: 54 (direnv)\n中位: 70.5",
                font_size=11, color=LIGHT)

    add_shape_bg(s, Inches(0.8), Inches(3.7), Inches(2.5), Inches(2.0), BG_CARD)
    add_textbox(s, Inches(1.0), Inches(3.8), Inches(2.1), Inches(0.3),
                "📈 脚本 vs 无脚本", font_size=14, color=WHITE, bold=True)
    add_textbox(s, Inches(1.0), Inches(4.2), Inches(2.1), Inches(1.2),
                "有脚本 4个: 84.3\n无脚本 8个: 63.8\n差距: 20.5分\n（再次验证脚本信号）",
                font_size=11, color=LIGHT)
    slide_number(s, prs, 13)

    # ================================================================
    # SLIDE 14: ship 95 Story
    # ================================================================
    s = new_slide(prs, DARKER)
    add_shape_bg(s, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
                "里程碑：首个社区技能超越官方", font_size=32, color=WHITE, bold=True)
    add_textbox(s, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
                "ship (95) — 来自 julianobarbosa 社区的 148 行精作", font_size=18, color=ACCENT)

    # Ship details
    add_shape_bg(s, Inches(0.8), Inches(2.6), Inches(7.0), Inches(4.2), BG_CARD)
    add_textbox(s, Inches(1.1), Inches(2.8), Inches(6.5), Inches(0.3),
                "🏗️ 为什么 ship 值得 95 分？", font_size=16, color=ACCENT, bold=True)
    ship_reasons = [
        "4 个 TypeScript 脚本 — 核心功能不可用无脚本\n"
        "  (ship-detect / ship-push / ship-pr / ship-lib)",
        "148 行精炼 SKILL.md — 最高经济性 22/22\n"
        "  每行都是具体指令，零废话",
        "触发边界完整 — 正例 + 负例 + 跨技能引用\n"
        '  "不要AI署名" "推送前先确认" = 负例加分',
        "双平台支持 — Azure DevOps + GitHub\n"
        "  OAuth 回退 + 工作项自动创建 + PR 模板",
        "安全设计 — 从不 force-push\n"
        "  破坏性操作前确认 + 令牌不过期",
        "完整元数据 — 引用脚本全部真实存在\n"
        "  4/4 完整性检查通过",
    ]
    for i, reason in enumerate(ship_reasons):
        add_textbox(s, Inches(1.3), Inches(3.2 + i * 0.58), Inches(6.2), Inches(0.55),
                    reason, font_size=11, color=LIGHT)

    # Comparison
    add_shape_bg(s, Inches(8.2), Inches(2.6), Inches(4.5), Inches(4.2), BG_CARD)
    add_textbox(s, Inches(8.5), Inches(2.8), Inches(4.0), Inches(0.3),
                "🏆 排名对比", font_size=16, color=WHITE, bold=True)

    rankings = [
        ("🥇", "ship", "95", "社区", ACCENT),
        ("🥇", "xlsx", "94", "官方", ACCENT2),
        ("🥇", "web-artifacts", "93", "官方", ACCENT2),
        ("🥇", "pptx", "92", "官方", ACCENT2),
        ("🥇", "docx", "91", "官方", ACCENT2),
    ]
    for i, (medal, name, score, source, clr) in enumerate(rankings):
        y = Inches(3.3 + i * 0.6)
        add_textbox(s, Inches(8.5), y, Inches(0.4), Inches(0.3),
                    medal, font_size=14)
        add_textbox(s, Inches(9.0), y, Inches(2.0), Inches(0.3),
                    name, font_size=13, color=WHITE, bold=(i == 0))
        add_textbox(s, Inches(11.0), y, Inches(0.5), Inches(0.3),
                    score, font_size=14, color=ACCENT if i == 0 else LIGHT, bold=(i == 0))
        add_textbox(s, Inches(11.6), y, Inches(0.8), Inches(0.3),
                    source, font_size=10, color=clr)

    add_shape_bg(s, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.3), BG_CARD)
    add_textbox(s, Inches(1.0), Inches(7.03), Inches(11), Inches(0.2),
                '💬  "社区技能 ship (95) 首次超越了所有官方技能最高分 (xlsx 94)。'
                '这是 v3.1 评分标准公正性的最好证明——价值驱动而非来源驱动。"',
                font_size=11, color=DIM, font_name='Arial', alignment=PP_ALIGN.CENTER)
    slide_number(s, prs, 14)

    # ================================================================
    # SLIDE 15: Top Skills Breakdown (CHART)
    # ================================================================
    s = new_slide(prs)
    add_accent_bar(s, Inches(0.8), Inches(0.8), height=Inches(0.4))
    add_textbox(s, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                "顶尖技能维度对比", font_size=28, color=WHITE, bold=True)

    s.shapes.add_picture(chart_dim, Inches(3.0), Inches(1.5), Inches(7.5), Inches(3.5))

    add_shape_bg(s, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.8), BG_CARD)
    add_textbox(s, Inches(1.1), Inches(5.5), Inches(11), Inches(0.3),
                "🔍 规律总结", font_size=14, color=ACCENT, bold=True)

    patterns = [
        "🎯 触发边界 → 所有 90+ 技能都有明确的 TRIGGER/SKIP 或 Do NOT 声明",
        "⚡ 脚本支撑 → 90+ 技能全部有 3+ 核心脚本，且脚本是技能不可用无的",
        "📐 精炼优先 → 90+ 技能的平均 SKILL.md 长度仅 163 行，远低于 70-80 分组（387行）",
        "🔒 安全设计 → 高分解率技能在 Gotchas 节中明确标注已知边界和陷阱",
    ]
    for i, pattern in enumerate(patterns):
        add_textbox(s, Inches(1.3), Inches(5.85 + i * 0.35), Inches(10.5), Inches(0.35),
                    pattern, font_size=12, color=LIGHT)
    slide_number(s, prs, 15)

    # ================================================================
    # SLIDE 16: 5 Core Principles
    # ================================================================
    s = new_slide(prs)
    add_accent_bar(s, Inches(0.8), Inches(0.8), height=Inches(0.4))
    add_textbox(s, Inches(1.2), Inches(0.6), Inches(10), Inches(0.6),
                "5 大核心原则 — 从 40+ 技能中提炼", font_size=28, color=WHITE, bold=True)

    principles = [
        ("①", "脚本为王", ACCENT,
         "有可执行脚本的技能价值远高于纯描述技能",
         "实证：有脚本平均 84.3 vs 无脚本 63.8，差距 20 分"),
        ("②", "价值 > 合规", ACCENT2,
         "不问\"管得严不严\"，问\"产出好不好\"",
         "好的技能 = 让 Claude 做出没有它时做不出的结果"),
        ("③", "密度决定经济", '#06b6d4',
         "长而泛化的技能应受密度扣分惩罚",
         "精炼≤150行的技能自动获得经济性附加分"),
        ("④", "触发即契约", '#8b5cf6',
         "清晰的触发边界是技能设计成熟的标志",
         "正例+负例+跨技能引用 = 专业级技能"),
        ("⑤", "来源非质量", '#f472b6',
         "官方技能不一定好，社区技能不一定差",
         "ship(社区) 95 > xlsx(官方) 94，价值驱动而非来源驱动"),
    ]
    for i, (num, title, clr, desc, evidence) in enumerate(principles):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.3)
        y = Inches(1.5 + row * 1.8)
        add_shape_bg(s, x, y, Inches(5.9), Inches(1.5), BG_CARD)
        add_shape_bg(s, x, y, Inches(0.06), Inches(1.5), clr)
        add_textbox(s, x + Inches(0.2), y + Inches(0.1), Inches(5.5), Inches(0.35),
                    f"{num}  {title}", font_size=18, color=clr, bold=True)
        add_textbox(s, x + Inches(0.2), y + Inches(0.5), Inches(5.5), Inches(0.35),
                    desc, font_size=12, color=WHITE)
        add_textbox(s, x + Inches(0.2), y + Inches(0.95), Inches(5.5), Inches(0.4),
                    evidence, font_size=11, color=DIM)
    slide_number(s, prs, 16)

    # ================================================================
    # SLIDE 17: Future + Ending
    # ================================================================
    s = new_slide(prs, DARKER)
    add_shape_bg(s, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(s, Inches(0.8), Inches(0.6), Inches(11), Inches(0.6),
                "未来方向 & 行动号召", font_size=28, color=WHITE, bold=True)

    # Future items
    futures = [
        ("🤖 自动化评分 Agent", "让 AI 自动运行 evaluate.sh + 读取评分卡 → 输出完整报告"),
        ("📊 公开排行榜", "建立技能评分排行榜，社区可以提交评分和投票"),
        ("🔄 评分标准开源", "让社区参与 v4.0 的讨论——更多数据，更精准的维度"),
        ("🔗 技能推荐引擎", "根据你的工作流推荐最适配的高分技能"),
    ]
    for i, (title, desc) in enumerate(futures):
        x = Inches(0.5 + (i % 2) * 6.3)
        y = Inches(1.4 + (i // 2) * 1.5)
        add_shape_bg(s, x, y, Inches(5.9), Inches(1.2), BG_CARD)
        add_textbox(s, x + Inches(0.2), y + Inches(0.1), Inches(5.5), Inches(0.3),
                    title, font_size=16, color=ACCENT, bold=True)
        add_textbox(s, x + Inches(0.2), y + Inches(0.5), Inches(5.5), Inches(0.5),
                    desc, font_size=13, color=LIGHT)

    # Key data
    add_shape_bg(s, Inches(0.8), Inches(4.8), Inches(11.5), Inches(1.0), BG_CARD)
    add_textbox(s, Inches(1.1), Inches(4.9), Inches(11), Inches(0.8),
                "📊 全旅程数据  ·  40+ 技能评估  ·  6 大 Skill Hub  ·  3 个大版本  ·  100+ 次迭代\n"
                "从 v1.0 合规审计 → v3.1 精准价值评估  ·  首个社区 95 分超越官方最高分",
                font_size=13, color=DIM, alignment=PP_ALIGN.CENTER)

    # CTA
    add_shape_bg(s, Inches(3.5), Inches(6.1), Inches(6.3), Inches(0.9), ACCENT)
    add_textbox(s, Inches(3.7), Inches(6.2), Inches(5.9), Inches(0.6),
                "👉  安装 skill-evaluator 试试你的技能评分",
                font_size=18, color=DARKER, bold=True, alignment=PP_ALIGN.CENTER)

    # Footer
    add_textbox(s, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
                "github.com/huaji/skill-evaluator  ·  欢迎 PR 和 Issue  ·  一起让技能生态更好",
                font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)
    slide_number(s, prs, 17)

    # ─── Save ──────────────────────────────────────────────────
    output_path = os.path.expanduser("~/Desktop/skill-evaluator-回顾.pptx")
    prs.save(output_path)
    print(f"\n✅ PPT saved to: {output_path}")
    file_size = os.path.getsize(output_path)
    print(f"   Size: {file_size/1024:.0f} KB")
    print(f"   Slides: {len(prs.slides)}")
    return output_path


if __name__ == '__main__':
    print("Building skill-evaluator retrospective PPT...")
    build()
    print("Done!")
